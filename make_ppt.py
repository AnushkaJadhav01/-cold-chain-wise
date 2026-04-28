from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

# Colors
DARK_BG   = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT    = RGBColor(0x00, 0xC2, 0xFF)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TXT = RGBColor(0xCC, 0xE5, 0xFF)
GREEN     = RGBColor(0x00, 0xE6, 0x76)
YELLOW    = RGBColor(0xFF, 0xD6, 0x00)
RED_CLR   = RGBColor(0xFF, 0x4D, 0x4D)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # completely blank

def add_slide():
    return prs.slides.add_slide(blank_layout)

def bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def txbox(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, wrap=True):
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf

def accent_bar(slide, t=0.85):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(t), Inches(13.33), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

def header_band(slide, title):
    band = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(0x05, 0x2A, 0x4A)
    band.line.fill.background()
    txbox(slide, title, 0.3, 0.15, 12, 0.8, size=26, bold=True, color=ACCENT)
    accent_bar(slide, t=1.1)

def bullets(slide, items, l, t, w, h, size=15, color=WHITE, gap=0.38):
    for i, item in enumerate(items):
        txbox(slide, f"  \u2022  {item}", l, t + i*gap, w, gap+0.05,
              size=size, color=color)

# ── Slide 1: Title ──────────────────────────────────────────────────
s = add_slide(); bg(s)
grad = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
grad.fill.solid(); grad.fill.fore_color.rgb = RGBColor(0x05, 0x14, 0x28)
grad.line.fill.background()
accent_bar(s, 3.4)
txbox(s, "GAT-RL COLD CHAIN INTELLIGENCE", 0.5, 1.0, 12.3, 1.0,
      size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txbox(s, "AI-Driven Multi-Depot Cold Chain Optimization", 0.5, 2.1, 12.3, 0.7,
      size=22, bold=False, color=LIGHT_TXT, align=PP_ALIGN.CENTER)
txbox(s, "Solution Challenge 2026  \u2022  Smart Supply Chains Track", 0.5, 2.9, 12.3, 0.5,
      size=16, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, "Team Leader: Anushka Anil Jadhav", 0.5, 5.8, 12.3, 0.5,
      size=14, color=LIGHT_TXT, align=PP_ALIGN.CENTER)
txbox(s, "Google Cloud Run  \u2022  Gemini 1.5 Flash  \u2022  GAT + RL", 0.5, 6.4, 12.3, 0.5,
      size=13, color=ACCENT, align=PP_ALIGN.CENTER)

# ── Slide 2: Team Details ───────────────────────────────────────────
s = add_slide(); bg(s); header_band(s, "Team Details")
rows = [
    ("Team Name",    "GAT-RL Cold Chain Intelligence"),
    ("Team Leader",  "Anushka Anil Jadhav"),
    ("Track",        "Smart Supply Chains – Resilient Logistics & Dynamic Optimization"),
    ("Cloud",        "Google Cloud Run"),
    ("Google AI",    "Gemini 1.5 Flash"),
]
for i,(k,v) in enumerate(rows):
    y = 1.4 + i*0.88
    box = s.shapes.add_shape(1, Inches(0.4), Inches(y), Inches(12.5), Inches(0.7))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x0A, 0x28, 0x45)
    box.line.color.rgb = ACCENT
    txbox(s, k+":", 0.5, y+0.08, 2.5, 0.55, size=14, bold=True, color=ACCENT)
    txbox(s, v,     3.1, y+0.08, 9.7, 0.55, size=14, color=WHITE)

# ── Slide 3: Problem Statement ──────────────────────────────────────
s = add_slide(); bg(s); header_band(s, "Problem Statement")
stat = s.shapes.add_shape(1, Inches(0.4), Inches(1.3), Inches(5.8), Inches(1.2))
stat.fill.solid(); stat.fill.fore_color.rgb = RED_CLR
stat.line.fill.background()
txbox(s, "\u20b955,000 Crore", 0.5, 1.35, 5.6, 0.7, size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, "Annual loss from cold-chain spoilage in India", 0.5, 2.05, 5.6, 0.4, size=12, color=WHITE, align=PP_ALIGN.CENTER)
probs = [
    "30-40% of perishable produce lost in transit annually",
    "Static routing ignores real-time disruptions (weather, traffic, temperature)",
    "Disruptions detected AFTER delivery timelines are already compromised",
    "~70% cold storage used for potatoes; pharma & dairy under-served",
    "No predictive AI to preemptively flag and re-route around bottlenecks",
    "Cascading delays impact farmers, retailers, consumers & pharma supply chains",
]
bullets(s, probs, 6.5, 1.3, 6.6, 5.5, size=14, gap=0.82)

# ── Slide 4: Solution Overview ─────────────────────────────────────
s = add_slide(); bg(s); header_band(s, "Brief About Our Solution")
txbox(s, "Cold-Chain-Wise: AI-Powered Supply Chain Resilience Platform", 0.4, 1.2, 12.5, 0.6,
      size=18, bold=True, color=ACCENT)
items = [
    "GAT-RL Hybrid AI: Graph Attention Networks + Reinforcement Learning for adaptive multi-depot routing",
    "What-If Simulator (Digital Twin): Real-time risk scoring via interactive sliders (delay, temp, distance)",
    "Google Gemini 1.5 Flash: Generates natural-language explanations for every risk scenario",
    "Autonomous AI Agent: Auto-triggers diagnostics when shipment risk > 85%; recommends emergency reroutes",
    "Real-Time Alert Feed: Live thermal breach & transit delay notifications",
    "Cloud-Native: Fully deployed on Google Cloud Run — serverless, scalable, zero infrastructure overhead",
]
bullets(s, items, 0.4, 1.9, 12.5, 5.2, size=14, gap=0.84)

# ── Slide 5: Opportunities / USP ───────────────────────────────────
s = add_slide(); bg(s); header_band(s, "Opportunities & USP")
txbox(s, "How it differs from existing solutions:", 0.4, 1.2, 12.5, 0.45, size=16, bold=True, color=ACCENT)
diff = [
    "Existing tools use static TSP/VRP; we use adaptive GAT-RL that learns from real-time data",
    "Competitors detect problems after spoilage; our system flags risk BEFORE transit is compromised",
    "First solution to combine GAT + RL + Google Gemini AI for explainable, actionable cold-chain insights",
]
bullets(s, diff, 0.4, 1.65, 12.5, 1.5, size=13, gap=0.5)
txbox(s, "Unique Selling Propositions (USP):", 0.4, 3.3, 12.5, 0.45, size=16, bold=True, color=GREEN)
usps = [
    "Digital Twin What-If Simulator — test hypothetical scenarios before dispatching shipments",
    "Gemini AI explains every risk score in plain language — no ML expertise needed by operators",
    "Fully cloud-native on Google Cloud Run — zero infra management, auto-scales with demand",
    "Targets India's Tier 2/3 agri belts & pharma hubs where cold-chain failure is most acute",
]
bullets(s, usps, 0.4, 3.78, 12.5, 3.2, size=13, color=LIGHT_TXT, gap=0.58)

# ── Slide 6: Features ──────────────────────────────────────────────
s = add_slide(); bg(s); header_band(s, "Key Features")
features = [
    ("What-If Simulator",         "Interactive sliders for delay/temp/distance → live color-coded risk score"),
    ("Gemini AI Explanations",    "Google Gemini 1.5 Flash generates risk reasoning & recommendations"),
    ("Live Telemetry Simulation", "Real-time tracking of temperature, humidity & transit progress"),
    ("AI Spoilage Prediction",    "Forecasts spoilage probability per shipment using predictive risk engine"),
    ("Autonomous AI Agent",       "Auto-triggers diagnostics > 85% risk; compressor health & reroute advice"),
    ("Dynamic Route Optimization","GAT-RL agent recommends optimized reroutes to nearest Tier-1 Cold Hubs"),
    ("Real-Time Alert Feed",      "Live notifications for thermal breaches, delays & critical events"),
    ("Performance Analytics",     "Before/after comparison charts, KPI tracking & benchmark results"),
    ("Delivery Planner",          "Multi-depot order planning with product constraints (vaccines, dairy, fruits)"),
    ("Firebase Auth",             "Secure login, saved plans & developer bypass mode for instant testing"),
]
for i,(name,desc) in enumerate(features):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.6
    y = 1.3 + row * 1.15
    box = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(6.3), Inches(1.0))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x08, 0x22, 0x3A)
    box.line.color.rgb = ACCENT
    txbox(s, name, x+0.1, y+0.05, 6.0, 0.4, size=13, bold=True, color=ACCENT)
    txbox(s, desc, x+0.1, y+0.48, 6.0, 0.45, size=11, color=LIGHT_TXT)

# ── Slide 7: Process Flow ──────────────────────────────────────────
s = add_slide(); bg(s); header_band(s, "Process Flow Diagram")
steps = [
    ("Logistics Manager\nInputs orders &\nconstraints", 0.3),
    ("GAT-RL Engine\nProcesses IoT,\ntraffic & weather", 2.9),
    ("Risk Scored &\nGemini Explanation\nGenerated", 5.5),
    ("Dynamic Route\nOptimization &\nFleet Dispatch", 8.1),
    ("Dashboard\nMonitoring &\nReal-Time Alerts", 10.7),
]
for label, x in steps:
    box = s.shapes.add_shape(9, Inches(x), Inches(2.2), Inches(2.3), Inches(2.5))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x05, 0x30, 0x55)
    box.line.color.rgb = ACCENT
    txbox(s, label, x+0.05, 2.55, 2.2, 1.8, size=12, bold=False,
          color=WHITE, align=PP_ALIGN.CENTER)
for x in [2.65, 5.25, 7.85, 10.45]:
    txbox(s, "\u2192", x, 3.05, 0.4, 0.5, size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
what_if = s.shapes.add_shape(1, Inches(4.5), Inches(5.1), Inches(4.3), Inches(0.85))
what_if.fill.solid(); what_if.fill.fore_color.rgb = RGBColor(0x00, 0x40, 0x60)
what_if.line.color.rgb = GREEN
txbox(s, "What-If Simulator: Operator tests scenarios before dispatch", 4.6, 5.2, 4.1, 0.6,
      size=12, color=GREEN, align=PP_ALIGN.CENTER)

# ── Slide 8: Architecture ──────────────────────────────────────────
s = add_slide(); bg(s); header_band(s, "Architecture Diagram")
layers = [
    (0.3, 1.3, 12.7, 1.0, RGBColor(0x04,0x20,0x38), ACCENT,    "FRONTEND LAYER  |  React 18 + TypeScript  •  Vite  •  Tailwind CSS  •  Shadcn UI  •  Framer Motion"),
    (0.3, 2.5, 12.7, 1.0, RGBColor(0x04,0x28,0x40), GREEN,     "BACKEND LAYER   |  FastAPI (Python 3.x)  •  GAT-RL Model  •  /gemini-explain Endpoint"),
    (0.3, 3.7, 12.7, 1.0, RGBColor(0x04,0x30,0x48), YELLOW,    "GOOGLE AI LAYER  |  Gemini 1.5 Flash API  •  Natural-Language Risk Explanations & Recommendations"),
    (0.3, 4.9, 12.7, 1.0, RGBColor(0x04,0x38,0x50), ACCENT,    "CLOUD LAYER       |  Google Cloud Run (Frontend + Backend)  •  Firebase Auth + Realtime DB"),
    (0.3, 6.1, 12.7, 0.85,RGBColor(0x04,0x18,0x2E), LIGHT_TXT, "DATA INPUTS         |  IoT Sensors (Temp, Humidity, GPS)  •  Traffic APIs  •  Weather Feeds"),
]
for (x,y,w,h,bg_c,border,txt) in layers:
    box = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = bg_c
    box.line.color.rgb = border
    txbox(s, txt, x+0.15, y+0.18, w-0.3, h-0.2, size=13, bold=False, color=WHITE)

# ── Slide 9: Tech Stack ─────────────────────────────────────────────
s = add_slide(); bg(s); header_band(s, "Technologies Used")
techs = [
    ("Google AI",        "Gemini 1.5 Flash",                 "Natural-language risk explanations"),
    ("Cloud",            "Google Cloud Run",                  "Serverless container deployment"),
    ("Frontend",         "React 18 + TypeScript",             "Component-based dashboard UI"),
    ("Styling",          "Tailwind CSS + Shadcn UI",          "Consistent design system"),
    ("Animations",       "Framer Motion",                     "Micro-animations & transitions"),
    ("Auth / DB",        "Firebase Auth + Realtime DB",       "Secure auth & data persistence"),
    ("Backend",          "FastAPI / Flask (Python 3.x)",      "API layer & GAT-RL logic"),
    ("AI Model",         "Graph Attention Networks (GAT)",    "Learns route importance from graphs"),
    ("AI Model",         "Reinforcement Learning (RL)",       "Adaptive real-time routing"),
    ("Container",        "Docker",                            "Consistent cloud deployment"),
]
headers = ["Layer", "Technology", "Purpose"]
col_w = [2.5, 4.2, 5.6]
col_x = [0.3, 2.85, 7.1]
for ci,(hdr,cw,cx) in enumerate(zip(headers,col_w,col_x)):
    hbox = s.shapes.add_shape(1, Inches(cx), Inches(1.2), Inches(cw), Inches(0.45))
    hbox.fill.solid(); hbox.fill.fore_color.rgb = ACCENT
    hbox.line.fill.background()
    txbox(s, hdr, cx+0.05, 1.25, cw-0.1, 0.38, size=13, bold=True, color=DARK_BG)
for ri, row in enumerate(techs):
    y = 1.7 + ri*0.52
    fill_c = RGBColor(0x08,0x22,0x3A) if ri%2==0 else RGBColor(0x06,0x1A,0x2E)
    for ci,(cell,cw,cx) in enumerate(zip(row,col_w,col_x)):
        box = s.shapes.add_shape(1, Inches(cx), Inches(y), Inches(cw), Inches(0.5))
        box.fill.solid(); box.fill.fore_color.rgb = fill_c
        box.line.color.rgb = RGBColor(0x10,0x40,0x60)
        col = ACCENT if ci==0 else (GREEN if ci==1 else WHITE)
        txbox(s, cell, cx+0.07, y+0.06, cw-0.1, 0.4, size=12, color=col)

# ── Slide 10: Results ──────────────────────────────────────────────
s = add_slide(); bg(s); header_band(s, "Results & Expected Impact")
metrics = [
    ("30-40%", "Spoilage\nReduction",    GREEN),
    ("85%→95%","On-Time\nDelivery",      ACCENT),
    ("\u20b91,500Cr","Monthly\nCost Savings", YELLOW),
    ("$105M+", "Annual\nROI",            RED_CLR),
]
for i,(val,label,col) in enumerate(metrics):
    x = 0.5 + i*3.2
    box = s.shapes.add_shape(9, Inches(x), Inches(1.4), Inches(2.8), Inches(2.5))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x06,0x1E,0x38)
    box.line.color.rgb = col
    txbox(s, val,   x+0.1, 1.55, 2.6, 0.9, size=30, bold=True,  color=col, align=PP_ALIGN.CENTER)
    txbox(s, label, x+0.1, 2.5,  2.6, 0.9, size=14, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, "Future Development Roadmap", 0.3, 4.1, 12.7, 0.45, size=16, bold=True, color=ACCENT)
future = [
    "Phase 2: Live GPS integration with fleet trucks",
    "Phase 3: Real IoT sensor data ingestion (temperature/humidity)",
    "Phase 4: Gemini-powered voice assistant for hands-free management",
    "Phase 5: Multi-language support (Hindi, Marathi, Tamil) for Tier 2/3 operators",
    "Phase 6: Industry-scale pilot with cold-chain logistics partners",
]
bullets(s, future, 0.3, 4.6, 12.7, 2.7, size=13, gap=0.5)

# ── Slide 11: Links ────────────────────────────────────────────────
s = add_slide(); bg(s); header_band(s, "Prototype Links")
links = [
    ("GitHub Repository",  "https://github.com/AnushkaJadhav01/cold-chain-wise.git",     ACCENT),
    ("Demo Video (3 min)", "https://youtu.be/ [paste your video link here]",              GREEN),
    ("Live MVP Link",      "https://[your-cloud-run-url].run.app",                         YELLOW),
    ("Working Prototype",  "https://[your-cloud-run-url].run.app",                         LIGHT_TXT),
]
for i,(label,url,col) in enumerate(links):
    y = 1.5 + i*1.35
    box = s.shapes.add_shape(1, Inches(0.4), Inches(y), Inches(12.5), Inches(1.1))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x06,0x1E,0x38)
    box.line.color.rgb = col
    txbox(s, f"{i+1}.  {label}", 0.55, y+0.08, 4.0, 0.45, size=15, bold=True, color=col)
    txbox(s, url, 4.7, y+0.12, 8.0, 0.4, size=12, color=WHITE)

# ── Slide 12: Thank You ────────────────────────────────────────────
s = add_slide(); bg(s)
accent_bar(s, 3.5)
txbox(s, "Thank You!", 0.5, 1.5, 12.3, 1.0, size=42, bold=True,
      color=ACCENT, align=PP_ALIGN.CENTER)
txbox(s, "Team GAT-RL Cold Chain Intelligence", 0.5, 2.8, 12.3, 0.6,
      size=20, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, "Leader: Anushka Anil Jadhav", 0.5, 3.5, 12.3, 0.5,
      size=16, color=LIGHT_TXT, align=PP_ALIGN.CENTER)
txbox(s, "Google Cloud Run  \u2022  Gemini 1.5 Flash  \u2022  GAT-RL  \u2022  Firebase", 0.5, 4.3, 12.3, 0.5,
      size=14, color=ACCENT, align=PP_ALIGN.CENTER)
txbox(s, "Solution Challenge 2026  |  Smart Supply Chains Track", 0.5, 6.5, 12.3, 0.5,
      size=12, color=LIGHT_TXT, align=PP_ALIGN.CENTER)

out = r"C:\Users\Admin\Desktop\Solution_Challenge_2026_GAT_RL.pptx"
prs.save(out)
print(f"Saved: {out}")
